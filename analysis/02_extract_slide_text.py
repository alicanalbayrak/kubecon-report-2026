"""Extract text from PDF and PPTX slide files."""

import json
import os
from multiprocessing import Process, Queue

import fitz  # PyMuPDF
from pptx import Presentation
from tqdm import tqdm

FILE_TIMEOUT = 30  # seconds per file

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

SOURCES = [
    {
        "path": os.path.join(REPO_ROOT, "kubecon-cloudnative-eu-2026", "kubecon_eu_2026_slides", "index.json"),
        "slides_dir": os.path.join(REPO_ROOT, "kubecon-cloudnative-eu-2026", "kubecon_eu_2026_slides"),
    },
    {
        "path": os.path.join(REPO_ROOT, "cncf-hosted-co-located-events-eu-2026", "colocated_eu_2026_slides", "index.json"),
        "slides_dir": os.path.join(REPO_ROOT, "cncf-hosted-co-located-events-eu-2026", "colocated_eu_2026_slides"),
    },
]

LOW_TEXT_THRESHOLD = 100


def collect_slide_paths() -> list[tuple[str, str]]:
    """Return list of (absolute_path, local_path) for all slide files."""
    paths = []
    seen = set()
    for source in SOURCES:
        with open(source["path"]) as f:
            data = json.load(f)
        for type_data in data["types"].values():
            for event in type_data["events"]:
                for file_info in event.get("files", []):
                    abs_path = os.path.join(source["slides_dir"], file_info["local_path"])
                    if abs_path not in seen and os.path.exists(abs_path):
                        seen.add(abs_path)
                        paths.append((abs_path, file_info["local_path"]))
    return paths


def _failed(error: str) -> dict:
    return {"text": "", "page_count": 0, "char_count": 0, "success": False, "error": error}


def extract_pdf(path: str) -> dict:
    try:
        doc = fitz.open(path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        text = "\n\n".join(pages)
        result = {
            "text": text,
            "page_count": len(doc),
            "char_count": len(text),
            "success": True,
        }
        doc.close()
        return result
    except Exception as e:
        return _failed(str(e))


def extract_pptx(path: str) -> dict:
    try:
        prs = Presentation(path)
        texts = []
        slide_count = 0
        for slide in prs.slides:
            slide_count += 1
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
        text = "\n".join(texts)
        return {
            "text": text,
            "page_count": slide_count,
            "char_count": len(text),
            "success": True,
        }
    except Exception as e:
        return _failed(str(e))


def _extract_worker(abs_path, ext, queue):
    """Worker function for subprocess-based extraction with timeout."""
    if ext == ".pdf":
        queue.put(extract_pdf(abs_path))
    else:
        queue.put(extract_pptx(abs_path))


def main():
    cache_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "cache")
    os.makedirs(cache_dir, exist_ok=True)

    paths = collect_slide_paths()
    print(f"Found {len(paths)} unique slide files")

    results = {}
    low_text = []

    for abs_path, local_path in tqdm(paths, desc="Extracting text"):
        ext = os.path.splitext(abs_path)[1].lower()
        if ext not in (".pdf", ".pptx"):
            continue

        # Run extraction in a subprocess so we can kill hung C code
        q = Queue()
        proc = Process(target=_extract_worker, args=(abs_path, ext, q))
        proc.start()
        proc.join(timeout=FILE_TIMEOUT)

        if proc.is_alive():
            proc.kill()
            proc.join()
            result = _failed("Timed out")
            tqdm.write(f"  TIMEOUT ({FILE_TIMEOUT}s): {local_path}")
        elif not q.empty():
            result = q.get()
        else:
            result = _failed("Worker crashed")

        results[local_path] = result
        if result["success"] and result["char_count"] < LOW_TEXT_THRESHOLD:
            low_text.append(local_path)

    success = sum(1 for r in results.values() if r["success"])
    failed = sum(1 for r in results.values() if not r["success"])
    print(f"Extracted: {success} success, {failed} failed")
    if low_text:
        print(f"Low text ({len(low_text)} files, <{LOW_TEXT_THRESHOLD} chars):")
        for p in low_text[:10]:
            print(f"  {p}")
        if len(low_text) > 10:
            print(f"  ... and {len(low_text) - 10} more")

    out_path = os.path.join(cache_dir, "slide_texts.json")
    with open(out_path, "w") as f:
        json.dump(results, f, indent=2)
    print(f"Written to {out_path}")


if __name__ == "__main__":
    main()
